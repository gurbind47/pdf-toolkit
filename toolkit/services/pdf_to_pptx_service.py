from io import BytesIO

import pypdfium2 as pdfium
from pptx import Presentation
from pptx.util import Emu

EMU_PER_POINT = 914400 / 72


def pdf_to_pptx(pdf_path, dpi=150):
    """Render each PDF page as an image on its own slide. Returns BytesIO."""
    pdf = pdfium.PdfDocument(pdf_path)
    try:
        if len(pdf) == 0:
            raise ValueError("PDF has no pages")

        prs = Presentation()
        first = pdf[0]
        prs.slide_width = Emu(int(first.get_width() * EMU_PER_POINT))
        prs.slide_height = Emu(int(first.get_height() * EMU_PER_POINT))
        blank_layout = prs.slide_layouts[6]

        scale = dpi / 72.0
        for i in range(len(pdf)):
            page = pdf[i]
            img = page.render(scale=scale).to_pil().convert("RGB")
            buf = BytesIO()
            img.save(buf, format="JPEG", quality=85)
            buf.seek(0)

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(buf, 0, 0, width=prs.slide_width, height=prs.slide_height)

        out = BytesIO()
        prs.save(out)
        out.seek(0)
        return out
    finally:
        pdf.close()
