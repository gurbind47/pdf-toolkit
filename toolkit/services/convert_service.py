import os
from io import BytesIO


def convert_to_pdf(file_path):
    """Convert a non-PDF file to PDF. Returns BytesIO."""
    ext = os.path.splitext(file_path)[1].lower()

    converters = {
        ".txt": _txt_to_pdf,
        ".md": _md_to_pdf,
        ".html": _html_to_pdf,
        ".htm": _html_to_pdf,
        ".docx": _docx_to_pdf,
        ".xlsx": _xlsx_to_pdf,
        ".pptx": _pptx_to_pdf,
    }

    converter = converters.get(ext)
    if not converter:
        raise ValueError(f"Unsupported format: {ext}")

    return converter(file_path)


def _txt_to_pdf(file_path):
    from fpdf import FPDF

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Courier", size=10)

    for line in text.split("\n"):
        pdf.cell(0, 5, line, new_x="LMARGIN", new_y="NEXT")

    buf = BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf


def _md_to_pdf(file_path):
    import markdown as md_lib
    from fpdf import FPDF

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        md_text = f.read()

    html = md_lib.markdown(md_text, extensions=["tables", "fenced_code"])

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    pdf.write_html(html)

    buf = BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf


def _html_to_pdf(file_path):
    import bleach
    from fpdf import FPDF

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        html = f.read()

    allowed_tags = [
        "h1", "h2", "h3", "h4", "h5", "h6", "p", "br", "hr",
        "b", "i", "u", "strong", "em", "a", "ul", "ol", "li",
        "table", "thead", "tbody", "tr", "th", "td",
        "blockquote", "pre", "code", "img", "sub", "sup",
    ]
    html = bleach.clean(html, tags=allowed_tags, attributes={"a": ["href"], "img": ["src", "alt"]})

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    pdf.write_html(html)

    buf = BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf


def _docx_to_pdf(file_path):
    from docx import Document
    from fpdf import FPDF

    doc = Document(file_path)

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    for para in doc.paragraphs:
        text = para.text
        if not text.strip():
            pdf.ln(3)
            continue

        style_name = (para.style.name or "").lower()

        if "heading 1" in style_name:
            pdf.set_font("Helvetica", "B", 20)
            pdf.cell(0, 10, text, new_x="LMARGIN", new_y="NEXT")
        elif "heading 2" in style_name:
            pdf.set_font("Helvetica", "B", 16)
            pdf.cell(0, 8, text, new_x="LMARGIN", new_y="NEXT")
        elif "heading 3" in style_name:
            pdf.set_font("Helvetica", "B", 13)
            pdf.cell(0, 7, text, new_x="LMARGIN", new_y="NEXT")
        elif "list" in style_name:
            pdf.set_font("Helvetica", "", 11)
            pdf.cell(10)
            pdf.cell(0, 6, f"  {text}", new_x="LMARGIN", new_y="NEXT")
        else:
            pdf.set_font("Helvetica", "", 11)
            pdf.multi_cell(0, 6, text)

        pdf.ln(1)

    # Tables
    for table in doc.tables:
        pdf.ln(3)
        pdf.set_font("Helvetica", "", 9)
        col_count = len(table.columns)
        col_w = (pdf.w - pdf.l_margin - pdf.r_margin) / max(col_count, 1)

        for row in table.rows:
            for cell in row.cells:
                pdf.cell(col_w, 6, cell.text[:50], border=1)
            pdf.ln()
        pdf.ln(3)

    buf = BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf


def _xlsx_to_pdf(file_path):
    from openpyxl import load_workbook
    from fpdf import FPDF

    wb = load_workbook(file_path, data_only=True)

    pdf = FPDF(orientation="L")
    pdf.set_auto_page_break(auto=True, margin=15)

    for sheet in wb.worksheets:
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 10, sheet.title, new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 8)

        max_col = sheet.max_column or 1
        col_w = (pdf.w - pdf.l_margin - pdf.r_margin) / max_col

        for row in sheet.iter_rows(values_only=True):
            for val in row:
                text = str(val) if val is not None else ""
                pdf.cell(col_w, 5, text[:40], border=1)
            pdf.ln()

    buf = BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf


def _pptx_to_pdf(file_path):
    from pptx import Presentation
    from fpdf import FPDF

    prs = Presentation(file_path)

    slide_w = prs.slide_width.inches if prs.slide_width else 10
    slide_h = prs.slide_height.inches if prs.slide_height else 7.5

    pdf = FPDF(orientation="L", unit="in", format=(slide_h, slide_w))

    for slide in prs.slides:
        pdf.add_page()
        pdf.set_font("Helvetica", "", 12)

        y_pos = 0.5
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    font_size = 12
                    style = ""
                    if para.runs:
                        run = para.runs[0]
                        if run.font.size:
                            font_size = min(run.font.size.pt, 36)
                        if run.font.bold:
                            style = "B"

                    pdf.set_font("Helvetica", style, font_size)
                    pdf.set_xy(0.5, y_pos)
                    pdf.multi_cell(slide_w - 1, font_size / 72 * 1.5, text)
                    y_pos = pdf.get_y() + 0.1

            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                col_count = len(table.columns)
                col_w = (slide_w - 1) / max(col_count, 1)
                pdf.set_font("Helvetica", "", 8)

                for row in table.rows:
                    pdf.set_x(0.5)
                    for cell in row.cells:
                        pdf.cell(col_w, 0.25, cell.text[:30], border=1)
                    pdf.ln()
                    y_pos = pdf.get_y() + 0.1

    buf = BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf
